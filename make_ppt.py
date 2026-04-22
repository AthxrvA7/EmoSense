from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "EmoSense"
    subtitle.text = "Emotion-Aware Chatbot Application\nProject Overview & Features"

    # Slide 2: Project Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Project Overview"
    tf = content.text_frame
    tf.text = "EmoSense is an intelligent, emotion-aware chatbot designed to detect and respond to user emotions."
    p = tf.add_paragraph()
    p.text = "It provides an empathetic conversational experience by adjusting its responses based on the detected emotional state of the user."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Key Components:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Machine Learning Emotion Model"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tkinter-based Modern UI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Context & Decision Engine"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Data Analysis & PDF Report Generation"
    p.level = 1

    # Slide 3: Core Features (1/2)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Core Features (1/2)"
    tf = content.text_frame
    tf.text = "Emotion & Context Detection"
    p = tf.add_paragraph()
    p.text = "Classifies user input into emotions: happy, sad, angry, anxious, neutral."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Calculates emotion intensity and detects context (e.g., work, social, exam)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Adaptive Chatbot Engine"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Dynamic decision engine chooses response types: empathy, validation, calming, or reinforcement."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Avoids repetition by tracking recent responses."
    p.level = 1

    # Slide 4: Core Features (2/2)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Core Features (2/2)"
    tf = content.text_frame
    tf.text = "Modern User Interface"
    p = tf.add_paragraph()
    p.text = "Clean, dark-themed Tkinter GUI."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Interactive chat bubbles resembling modern messaging apps."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dashboard & Reporting"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "In-app dashboard displays emotional metrics, patterns, and risk assessments."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Plots emotion distributions and trends."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "One-click PDF report generation for a comprehensive emotional overview."
    p.level = 1

    # Slide 5: System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "System Architecture"
    tf = content.text_frame
    tf.text = "AuthSystem (auth.py)"
    p = tf.add_paragraph()
    p.text = "Manages user registration, login, and sessions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "ChatbotEngine (chatbot_engine.py)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Orchestrates emotion detection, context parsing, and response generation."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database (database.py)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Stores chat history securely for ongoing analysis."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "EmotionAnalyzer & Generator (analysis.py / report_generator.py)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Aggregates data to produce insights and downloadable PDF reports."
    p.level = 1

    # Slide 6: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusion"
    tf = content.text_frame
    tf.text = "Value Proposition:"
    p = tf.add_paragraph()
    p.text = "EmoSense bridges the gap between simple text bots and true empathetic interaction."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Enhancements:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Voice emotion recognition."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Advanced NLP models for nuanced emotion tracking."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Cloud synchronization across multiple devices."
    p.level = 1

    prs.save('EmoSense_Presentation.pptx')
    print("Presentation saved as EmoSense_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
